from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()

def add_title_slide(title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

def add_slide(title_text, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.level = 0
        if point.startswith("- ") or point.startswith("* "):
            p.level = 1

# Title
add_title_slide("SpotBot: AOI System for PCBs", "Deep Learning for Computer Vision Final Project")

# 1. Problem Statement
add_slide("1. Problem Statement", [
    "Manufacturing Printed Circuit Boards (PCBs) is highly complex.",
    "Even microscopic errors like missing holes, open circuits, or bridged traces can compromise entire electronic devices.",
    "Manual Visual Inspection (MVI) is slow, subjective, and prone to human fatigue.",
    "Defective boards often pass quality control, causing extremely costly failures in the field."
])

# 2. Abstract
add_slide("2. Abstract", [
    "SpotBot is a real-time Automated Optical Inspection (AOI) system designed to instantly detect manufacturing defects.",
    "Leverages Deep Learning (YOLOv8) for object detection.",
    "Utilizes traditional Computer Vision (OpenCV) for structural analysis and simulated thermal mapping.",
    "Identifies six specific defect classes instantly via a live, reactive web interface."
])

# 3. Methodology
add_slide("3. Methodology", [
    "Data Collection: PKU-Market-PCB dataset.",
    "Preprocessing: Image resizing, tensor normalization, and augmentation.",
    "Deep Learning: Training YOLOv8 for precise bounding box regression and classification.",
    "Computer Vision: Edge detection, HSV color filtering, adaptive thresholding via OpenCV.",
    "Web Architecture: FastAPI & React to stream asynchronous video chunks and render SPC Control Charts."
])

# 4. Existing Solution
add_slide("4. Existing Solutions", [
    "Human Operators: Use magnifying glasses. Very slow and prone to missed defects.",
    "Proprietary AOI Machines: Cost $50,000+ and require manual, rigid algorithmic programming for EVERY new board layout.",
    "Basic Machine Vision: Relies on pure template matching, which breaks immediately under different lighting or slight hardware rotations."
])

# 5. Proposed Solution
add_slide("5. Proposed Solution", [
    "An accessible, highly adaptive Deep Learning AOI system.",
    "Ingests live images or web camera feeds.",
    "Executes a custom-trained YOLOv8 model to pinpoint microscopic anomalies natively—regardless of lighting or slight rotations.",
    "Integrates mathematical computer vision to generate thermal damage logs."
])

# 6. Base Idea
add_slide("6. Base Idea", [
    "Replace rigid, inflexible algorithmic matching with an intelligent, generalized neural network.",
    "The AI inherently 'understands' what a broken hardware trace looks like.",
    "Augment this Deep AI with classical Computer Vision techniques to rip out visual features.",
    "Provide contextual views (architectural blueprints and glowing heatmaps) explicitly designed for human repair technicians."
])

# 7. Tech Stack: Basic vs. Our Idea
add_slide("7. Tech Stack: Basic vs. Our Idea", [
    "YOLOv8 Basic: Draws boxes around standard objects (cars/cats).",
    "YOLOv8 Ours: Maps precise coordinates around 6 microscopic PCB defects in milliseconds.",
    "OpenCV Basic: Reads and scripts image matrices.",
    "OpenCV Ours: Simulates thermal heatmaps via HSV filtering; generates architectural blueprints.",
    "Scikit-Learn Basic: Data science algorithms.",
    "Scikit-Learn Ours: Isolation Forests flag statistically catastrophic damages not trained in YOLO.",
    "FastAPI & React: Pipes continuous frame chunks from hardware into PyTorch matrices dynamically."
])

# 8. Features
add_slide("8. Features", [
    "Live Real-Time Hardware Detection (Webcam)",
    "URL Link & Clipboard Scraping Engine",
    "Simulated Thermal & HSV Heatmaps",
    "Blueprint & Wireframe PCB Rendering",
    "Live Statistical Process Control (SPC) Dashboard",
    "LLM Repair Guide Generation (Google Gemini)"
])

# 9. Modules
add_slide("9. Modules", [
    "Vision & Perception Module: YOLOv8 execution and OpenCV thresholding/filters.",
    "Backend API Module: FastAPI routing, UUID tracking, and WAF-bypassing Web Scrapers.",
    "Analytics & DB Module: MySQL relational storage and Scikit-Learn Isolation Forests for global anomaly scoring.",
    "User Interface Module: React DOM virtual updates, Recharts charts, and WebRTC Camera links."
])

# 10. Advantages
add_slide("10. Advantages", [
    "Lightning Fast: Inferences are executed in milliseconds.",
    "Highly Adaptable: Deep learning natively handles varying factory lighting conditions without breaking.",
    "Multi-Modal: Combines spatial coordinates with pure visual representations and LLM-generated text.",
    "Cost-Effective: Operates on standard digital cameras instead of heavily proprietary $50k industrial rigs."
])

# 11. Disadvantages
add_slide("11. Disadvantages", [
    "Hardware Dependent: Requires a decent CPU/GPU hardware accelerator to maintain a 30fps video stream without lagging.",
    "Data Hungry: Training for entirely new proprietary boards can require thousands of annotated manual images.",
    "Unseen Defect Limitations: Neural networks can occasionally struggle confidently identifying a completely new, untrained class of defect."
])

# 12. Justification
add_slide("12. Project Justification", [
    "Perfect synthesis of 'Deep Learning for Computer Vision'.",
    "Proves mastery over Deep Convolutional Neural Networks (training YOLO for spatial detection).",
    "Bridges perfectly with classical Computer Vision math (OpenCV matrix manipulation, thresholding, color space conversions).",
    "Proves we can apply theoretical algorithms directly to a physical, high-speed industrial robotics problem."
])

# 13. Conclusion
add_slide("13. Conclusion", [
    "SpotBot successfully demonstrates how state-of-the-art Deep Learning can drastically modernize hardware manufacturing.",
    "By fusing YOLOv8, OpenCV, and React, we delivered a full-stack Quality Assurance pipeline.",
    "Capable of outperforming traditional human-in-the-loop manual inspection."
])

# 14. Reference
add_slide("14. References", [
    "Ultralytics YOLOv8 Academic Documentation.",
    "PKU-Market-PCB Dataset (Intelligent Robotics Lab, Peking University).",
    "OpenCV Documentation (Edge Detection, Thresholds, Color Maps).",
    "Scikit-Learn Isolation Forest Literature.",
    "React.js and FastAPI software engineering documentation."
])

output_file = "SpotBot_Final_Presentation.pptx"
prs.save(output_file)
print(f"Presentation saved successfully as {output_file}")
